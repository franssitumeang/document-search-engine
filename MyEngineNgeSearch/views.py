from django.shortcuts import render
import PyPDF2
import docx2txt
from pptx import Presentation
from os import listdir
import os
from collections import defaultdict
import math
import string
import re
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from django.http import HttpResponse
import mimetypes
import urllib.parse
from django.core.paginator import Paginator
import time


#

class Document():
    def __init__(self, id, name, type, text):
        self.id = id
        self.name = name
        self.type = type
        self.text = text + ' '+name

#

extensions = ['pdf','docx','pptx']
def get_content_file(directory):
    document_contents = []
    id_doc = 0
    for ex in extensions:
        filenames = [f for f in listdir(directory) if f.endswith('.'+ex)]
        for filename in filenames:
            try:
                path = directory+'\\'+filename
                text = ''
                type = ''
                if(ex == 'pdf'):
                    type = 'pdf'
                    pdf_file = open(path,'rb')
                    pdf_reader = PyPDF2.PdfFileReader(pdf_file)
                    num_pages = pdf_reader.numPages
                    for i in range(num_pages):
                        page = pdf_reader.getPage(i)
                        text += page.extractText()
                elif(ex == 'docx'):
                    type = 'docx'
                    text = docx2txt.process(path)
                elif(ex == 'pptx'):
                    type = 'pptx'
                    ppt_file = Presentation(path)
                    for slide in ppt_file.slides:
                        for shape in slide.shapes:
                            if(shape.has_text_frame):
                                text += shape.text
                document = Document(id_doc, filename, type, text)
                document_contents.append(document)
                id_doc +=1
            except:
                pass
    return document_contents

#

all_document = get_content_file('D:\Kuliah\Jupyter Notebook\INRE\My Engine Nge-Search\Document Collection')
N = len(all_document)
dictionary = set()
postings = defaultdict(dict)
document_frequency = defaultdict(int)
length = defaultdict(float)

#

def tokenize(document):
    document = document.lower()
    document = re.sub(r'^https?:\/\/.*[\r\n]*', '', document, flags=re.MULTILINE)
    tokens = document.split()
    punc = string.punctuation
    tokens = [token.strip(punc) for token in tokens]

    stop_words = set(stopwords.words('english'))
    tokens = [w for w in tokens if not w in stop_words]

    stemmer = PorterStemmer()
    for i in range(0, len(tokens)):
        if (tokens[i] != stemmer.stem(tokens[i])):
            tokens[i] = stemmer.stem(tokens[i])
    return tokens

#

def set_terms_and_postings():
    global dictionary, postings
    for doc in all_document:
        terms = tokenize(doc.text)
        unique_terms = set(terms)
        dictionary = dictionary.union(unique_terms)
        for term in unique_terms:
            postings[term][doc.id] = terms.count(term)
set_terms_and_postings()

#

def set_document_frequencies():
    global document_frequency
    for term in dictionary:
        document_frequency[term] = len(postings[term])
set_document_frequencies()

#

def set_inverse_document_frequency(term):
    if term in dictionary:
        return math.log10(float(N)/float(document_frequency[term]))
    else:
        return 0.0

#

def imp(term,id):
    if id in postings[term]:
        return postings[term][id]*set_inverse_document_frequency(term)
    else:
        return 0.0

#

def set_lengths():
    global length
    for doc in all_document:
        l = 0
        for term in dictionary:
            l += imp(term,doc.id)**2
        length[doc.id] = math.sqrt(l)
set_lengths()

#

def similarity(query, id):
    similarity = 0.0
    for term in query:
        if term in dictionary:
            similarity += set_inverse_document_frequency(term) * imp(term, id)

    similarity = similarity / length[id]

    return similarity

#

def find_doc(id, extension = None):
    for d in all_document:
        if(extension):
            if(d.id == id and d.type == extension):
                return d.name
        else:
            if(d.id == id):
                return d.name
    return None

#

class ResultDocoment():
    def __init__(self, score, name):
        self.score = score
        self.name = name

#
def search(query):
    query = tokenize(query)
    query_extension = None
    if (query[len(query)-1] in extensions):
        query_extension = query[len(query)-1]
        query = query[:-1]
    all_result_docoment = []
    id_set = []
    for term in query:
        for id_d in postings[term].keys():
            id_set.append(id_d)
    id_set = set(id_set)
    m = ''
    if not id_set:
        m = "No documents matched all query terms."
    else:
        scores = sorted([(id,similarity(query,id))
                         for id in id_set],
                        key=lambda x: x[1],
                        reverse=True)
        for (id,score) in scores:
            if(query_extension):
                if(find_doc(id,extension=query_extension)):
                    all_result_docoment.append(ResultDocoment(score,find_doc(id,extension=query_extension)))
            else:
                all_result_docoment.append(ResultDocoment(score,find_doc(id)))
    return all_result_docoment,m



def index(request):
    content = {

    }
    return render(request, 'index.html', content)


def result(request):
    q = request.GET.get('q')
    page = request.GET.get('page')
    start_time = time.time()
    document,m = search(q)
    total_document = len(document)
    elapsed_time = time.time() - start_time
    for d in document:
        print(d.score, d.name)
    paginator = Paginator(document, 10)

    document = paginator.get_page(page)
    content = {
        'documents': document,
        'query': q,
        'm': m,
        'elapsed_time':elapsed_time,
        'total_document': total_document
    }
    return render(request, 'result.html', content)


def download(request, doc_name):
    filename = "D:/Kuliah/Jupyter Notebook/INRE/My Engine Nge-Search/Document Collection/"+doc_name
    fp = open(filename, 'rb')
    response = HttpResponse(fp.read())
    fp.close()
    type, encoding = mimetypes.guess_type(doc_name)
    if type is None:
        type = 'application/octet-stream'
    response['Content-Type'] = type
    response['Content-Length'] = str(os.stat(filename).st_size)
    if encoding is not None:
        response['Content-Encoding'] = encoding

    if u'WebKit' in request.META['HTTP_USER_AGENT']:
        filename_header = 'filename=%s' % doc_name.encode('utf-8')
    elif u'MSIE' in request.META['HTTP_USER_AGENT']:
        filename_header = ''
    else:
        filename_header = 'filename*=UTF-8\'\'%s' % urllib.parse.quote(doc_name.encode('utf-8'))
    response['Content-Disposition'] = 'attachment; ' + filename_header
    return response

def pagination(request):
    numbers = range(0,100)
    paginator = Paginator(numbers,10)
    page = request.GET.get('page')
    numbers = paginator.get_page(page)
    content = {
        'numbers' : numbers,
    }
    return render(request, 'pagination.html', content)
